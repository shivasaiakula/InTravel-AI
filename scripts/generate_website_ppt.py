from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path


TITLE_COLOR = RGBColor(0x0E, 0x2A, 0x47)
ACCENT_COLOR = RGBColor(0x0B, 0x74, 0xDE)
TEXT_COLOR = RGBColor(0x2F, 0x38, 0x4D)


def set_title_style(shape):
    tf = shape.text_frame
    for p in tf.paragraphs:
        if not p.runs:
            continue
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = TITLE_COLOR


def set_subtitle_style(shape):
    tf = shape.text_frame
    for p in tf.paragraphs:
        if not p.runs:
            continue
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(20)
        run.font.color.rgb = ACCENT_COLOR


def set_body_style(text_frame):
    for i, p in enumerate(text_frame.paragraphs):
        if not p.runs:
            continue
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(22 if i == 0 else 20)
            run.font.color.rgb = TEXT_COLOR


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    set_title_style(slide.shapes.title)
    set_subtitle_style(slide.placeholders[1])


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    set_title_style(slide.shapes.title)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, item in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = item
        p.level = 0

    set_body_style(body)


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = title
    set_title_style(title_shape)

    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.3))
    right_box = slide.shapes.add_textbox(Inches(6.2), Inches(1.5), Inches(6.7), Inches(5.3))

    left_tf = left_box.text_frame
    left_tf.text = left_title
    for item in left_items:
        p = left_tf.add_paragraph()
        p.text = item
        p.level = 1

    right_tf = right_box.text_frame
    right_tf.text = right_title
    for item in right_items:
        p = right_tf.add_paragraph()
        p.text = item
        p.level = 1

    set_body_style(left_tf)
    set_body_style(right_tf)


def add_timeline_slide(prs, title, points):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    set_title_style(slide.shapes.title)

    y = 1.5
    for phase, desc in points:
        line = slide.shapes.add_shape(1, Inches(0.8), Inches(y + 0.15), Inches(0.22), Inches(0.22))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_COLOR
        line.line.color.rgb = ACCENT_COLOR

        box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(11.5), Inches(0.6))
        tf = box.text_frame
        tf.text = f"{phase}: {desc}"
        set_body_style(tf)
        y += 0.95


def build_presentation(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "InTravel AI Website",
        "Product Overview and Booking Experience Presentation",
    )

    add_bullets_slide(
        prs,
        "Vision and Problem",
        [
            "Travel planning is fragmented across discovery, booking, and itinerary tools.",
            "InTravel AI unifies planning, hotel booking, transport booking, and post-booking actions.",
            "Goal: faster decisions, fewer drop-offs, and a trustworthy end-to-end booking flow.",
        ],
    )

    add_two_column_slide(
        prs,
        "Core Website Modules",
        "Frontend Experience",
        [
            "Home, Explore, Planner, Dashboard, Budget, Packing, Transport, Bookings",
            "Booking-first UX with hotel and transport flows",
            "Ticket actions: print, PDF, email, and calendar export",
        ],
        "Backend Services",
        [
            "Express API with booking, transport, and hotel recommendation routes",
            "DB-first with graceful local-memory fallback",
            "Payment status updates and ticket communication endpoints",
        ],
    )

    add_bullets_slide(
        prs,
        "Hotel Booking Excellence",
        [
            "Best hotel recommendations by city with ranking by rating and value.",
            "Case-insensitive and partial city matching for robust search.",
            "Live pricing with nights x rooms and suggested totals.",
            "Strong validations: date sanity, contact quality, and payable amount checks.",
        ],
    )

    add_bullets_slide(
        prs,
        "Transport Booking Excellence",
        [
            "RedBus-style service selection, filters, and realistic seat layout.",
            "Boarding and dropping selection with passenger and contact capture.",
            "Fallback transport options when live feeds are unavailable.",
            "Consistent checkout and ticketing experience with hotel flow parity.",
        ],
    )

    add_bullets_slide(
        prs,
        "Checkout and Payments",
        [
            "Fare breakup: base amount, taxes, service fee, discount, final payable.",
            "Coupon validation via API with type-aware offer rules.",
            "Payment lifecycle: pending, success, failed, and retry support.",
            "Payment state reflected in ticket summary and booking history badges.",
        ],
    )

    add_bullets_slide(
        prs,
        "Post-Booking and Reliability",
        [
            "Cancellation flow with refund estimation before confirmation.",
            "Ticket actions: PDF, print, QR rendering, and calendar export.",
            "Email ticket endpoint with SMTP and safe fallback mode.",
            "Quick booking modify endpoint for post-booking updates.",
        ],
    )

    add_two_column_slide(
        prs,
        "Technical Stack",
        "Frontend",
        [
            "React + Vite",
            "Axios for API interactions",
            "Modular CSS and responsive design",
            "Lucide icons for visual consistency",
        ],
        "Backend",
        [
            "Node.js + Express",
            "MySQL integration with fallback resilience",
            "REST APIs for hotels, transport, bookings, offers",
            "Nodemailer for ticket communication",
        ],
    )

    add_timeline_slide(
        prs,
        "Recommended Next Roadmap",
        [
            ("Phase 1", "Integrate real payment gateway and webhook verification"),
            ("Phase 2", "Add full modify-booking modal with repricing and differential payment"),
            ("Phase 3", "Improve recommendation personalization using traveler intent"),
            ("Phase 4", "Automate test coverage for booking and cancellation flows"),
            ("Phase 5", "Add analytics dashboard for conversion and refund metrics"),
        ],
    )

    add_bullets_slide(
        prs,
        "Business Value",
        [
            "Higher booking conversion through reduced friction and transparent pricing.",
            "Stronger trust via refund clarity, clear ticketing, and reliable backups.",
            "Improved retention with complete planning-to-booking lifecycle in one place.",
            "Foundation ready for scale with modular APIs and resilient architecture.",
        ],
    )

    prs.save(output_path)


if __name__ == "__main__":
    repo_root = Path(__file__).resolve().parents[1]
    output_path = repo_root / "InTravel-AI-Website-Overview.pptx"
    build_presentation(str(output_path))
